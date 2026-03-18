import asyncio
import logging
import mimetypes
import urllib.parse
from pathlib import Path
from typing import Any, Dict, Optional, Tuple

from fastapi import APIRouter, Depends, File, Form, HTTPException, Query, UploadFile
from fastapi.responses import FileResponse, HTMLResponse, StreamingResponse
from pptx import Presentation
from sqlalchemy.exc import IntegrityError
from sqlalchemy.orm import Session

from ...core.tools.adapters.vibe.file_tool import read_file
from ..auth_dependencies import get_current_user
from ..config import MAX_FILE_SIZE, UPLOADS_DIR, get_upload_path, is_allowed_file
from ..models.database import get_db
from ..models.uploaded_file import UploadedFile
from ..models.user import User
from .legacy_file import (
    infer_user_id_from_legacy_path,
    is_valid_uuid,
    resolve_legacy_file_path,
    resolve_legacy_file_path_cross_user,
)

logger = logging.getLogger(__name__)

file_router = APIRouter(prefix="/api/files", tags=["files"])


def _user_id_value(user: User) -> int:
    return int(getattr(user, "id"))


def _file_user_id_value(file_record: UploadedFile) -> int:
    return int(getattr(file_record, "user_id"))


def _is_admin_user(user: User) -> bool:
    return bool(getattr(user, "is_admin", False))


def _file_storage_path_value(file_record: UploadedFile) -> str:
    return str(getattr(file_record, "storage_path"))


def _file_name_value(file_record: UploadedFile) -> str:
    return str(getattr(file_record, "filename"))


def _parse_task_id(task_id: Optional[str]) -> Optional[int]:
    if task_id is None or task_id == "":
        return None
    try:
        return int(task_id)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail="Invalid task_id") from exc


def _guess_media_type(filename: str) -> str:
    media_type, _ = mimetypes.guess_type(filename)
    return media_type or "application/octet-stream"


def _build_unique_file_path(path: Path) -> Path:
    if not path.exists():
        return path
    stem = path.stem
    suffix = path.suffix
    parent = path.parent
    i = 1
    while True:
        candidate = parent / f"{stem}_{i}{suffix}"
        if not candidate.exists():
            return candidate
        i += 1


def _ensure_under_uploads(path: Path, user_id: int) -> None:
    resolved_path = path.resolve()
    uploads_root = UPLOADS_DIR.resolve()
    user_root = (UPLOADS_DIR / f"user_{user_id}").resolve()
    try:
        resolved_path.relative_to(uploads_root)
        resolved_path.relative_to(user_root)
    except ValueError as exc:
        raise HTTPException(status_code=403, detail="Access denied") from exc


def _resolve_public_preview_target(
    base_path: Path, relative_path: Optional[str], user_id: int
) -> Path:
    _ensure_under_uploads(base_path, user_id)
    if not relative_path:
        return base_path

    base_dir = base_path.parent.resolve()
    candidate = (base_dir / relative_path).resolve()

    try:
        candidate.relative_to(base_dir)
    except ValueError as exc:
        raise HTTPException(status_code=403, detail="Access denied") from exc

    _ensure_under_uploads(candidate, user_id)
    return candidate


def _to_unix_timestamp(path: Path, fallback: Any) -> int:
    if path.exists():
        return int(path.stat().st_mtime)
    if fallback is not None and hasattr(fallback, "timestamp"):
        return int(fallback.timestamp())
    return 0


def _extract_relative_path(storage_path: Path, user_id: int) -> str:
    user_root = UPLOADS_DIR / f"user_{user_id}"
    try:
        return str(storage_path.relative_to(user_root))
    except ValueError:
        return storage_path.name


def _collect_backfill_user_ids(user: User) -> list[int]:
    if not _is_admin_user(user):
        return [_user_id_value(user)]

    user_ids: list[int] = []
    if not UPLOADS_DIR.exists():
        return user_ids

    for child in UPLOADS_DIR.iterdir():
        if not child.is_dir() or not child.name.startswith("user_"):
            continue
        try:
            user_ids.append(int(child.name.replace("user_", "", 1)))
        except ValueError:
            continue
    return user_ids


def _infer_backfill_task_id(
    db: Session, file_path: Path, user_id: int
) -> Optional[int]:
    from ..models.task import Task

    user_root = UPLOADS_DIR / f"user_{user_id}"
    try:
        rel_parts = file_path.relative_to(user_root).parts
    except ValueError:
        return None

    if not rel_parts:
        return None
    first_part = rel_parts[0]
    task_id_part: Optional[str] = None
    if first_part.startswith("web_task_"):
        task_id_part = first_part.replace("web_task_", "", 1)
    elif first_part.startswith("task_"):
        task_id_part = first_part.replace("task_", "", 1)

    if task_id_part is None:
        return None

    try:
        task_id = int(task_id_part)
    except ValueError:
        return None

    task = db.query(Task.id).filter(Task.id == task_id, Task.user_id == user_id).first()
    return task_id if task is not None else None


def _backfill_uploaded_file_records(db: Session, user: User) -> None:
    if not UPLOADS_DIR.exists():
        return

    target_user_ids = _collect_backfill_user_ids(user)
    if not target_user_ids:
        return

    existing_paths = {
        row[0]
        for row in db.query(UploadedFile.storage_path)
        .filter(UploadedFile.user_id.in_(target_user_ids))
        .all()
    }

    created = 0
    for target_user_id in target_user_ids:
        user_root = UPLOADS_DIR / f"user_{target_user_id}"
        if not user_root.exists() or not user_root.is_dir():
            continue

        for candidate in user_root.rglob("*"):
            if not candidate.is_file():
                continue

            storage_path = str(candidate)
            if storage_path in existing_paths:
                continue

            file_record = UploadedFile(
                user_id=target_user_id,
                task_id=_infer_backfill_task_id(db, candidate, target_user_id),
                filename=candidate.name,
                storage_path=storage_path,
                mime_type=_guess_media_type(candidate.name),
                file_size=candidate.stat().st_size,
            )
            db.add(file_record)
            existing_paths.add(storage_path)
            created += 1

    if created > 0:
        try:
            db.commit()
            logger.info(f"Backfilled {created} uploaded_files records")
        except IntegrityError:
            db.rollback()
            logger.warning(
                "Backfill commit hit unique constraint race; rolled back safely"
            )


def _get_file_record(db: Session, file_id: str) -> UploadedFile:
    file_record = db.query(UploadedFile).filter(UploadedFile.file_id == file_id).first()
    if file_record is None:
        raise HTTPException(status_code=404, detail="File not found")
    return file_record


def _resolve_file_path(
    db: Session, file_id_or_path: str, user_id: int
) -> Tuple[Optional[UploadedFile], Path, int]:
    """
    Resolve file_id or legacy path to file record and actual path.

    This function handles both:
    - New system: UUID file_id that maps to a database record
    - Legacy system: Relative paths like "web_task_235/output/file.jpeg"
    - Workspace system: Files created by agents in workspace directories

    Args:
        db: Database session
        file_id_or_path: Either a UUID file_id or a legacy file path
        user_id: Current user's ID for permission checks

    Returns:
        Tuple of (file_record or None, file_path, owner_user_id)

    Raises:
        HTTPException: If file is not found
    """
    # If it's a valid UUID, try to find by file_id (new system)
    if is_valid_uuid(file_id_or_path):
        file_record = (
            db.query(UploadedFile)
            .filter(UploadedFile.file_id == file_id_or_path)
            .first()
        )
        if file_record:
            return (
                file_record,
                Path(_file_storage_path_value(file_record)),
                _file_user_id_value(file_record),
            )

    # For legacy paths, resolve from filesystem
    # First try to find in current user's directory
    file_path = resolve_legacy_file_path(file_id_or_path, user_id)
    owner_user_id = user_id

    # If not found and user is admin, try to infer owner from path and search all users
    if file_path is None:
        # Try to infer the correct user_id from the path
        inferred_user_id = infer_user_id_from_legacy_path(db, file_id_or_path)
        if inferred_user_id is not None:
            file_path = resolve_legacy_file_path(file_id_or_path, inferred_user_id)
            if file_path is not None:
                owner_user_id = inferred_user_id

        # If still not found and user is admin, try searching in all user directories
        if file_path is None and _is_admin_user_by_id(db, user_id):
            result = resolve_legacy_file_path_cross_user(file_id_or_path)
            if result is not None:
                file_path, owner_user_id = result

    if file_path is None:
        raise HTTPException(status_code=404, detail="File not found")

    # Try to find a matching database record (might exist for backfilled files)
    file_record = (
        db.query(UploadedFile)
        .filter(UploadedFile.storage_path == str(file_path))
        .first()
    )

    return (file_record, file_path, owner_user_id)


def _is_admin_user_by_id(db: Session, user_id: int) -> bool:
    """Check if a user is admin by user ID."""
    from ..models.user import User

    user = db.query(User).filter(User.id == user_id).first()
    return user is not None and getattr(user, "is_admin", False)


def _check_file_access(file_record: UploadedFile, user: User) -> None:
    if _is_admin_user(user):
        return
    if _file_user_id_value(file_record) != _user_id_value(user):
        raise HTTPException(status_code=403, detail="Access denied")


async def _try_convert_pptx_to_pdf(path: Path) -> Optional[StreamingResponse]:
    if path.suffix.lower() != ".pptx":
        return None

    # Check for cached PDF next to the PPTX file
    cached_pdf = path.with_suffix(".pdf")
    if cached_pdf.exists():
        pptx_mtime = path.stat().st_mtime
        pdf_mtime = cached_pdf.stat().st_mtime
        if pdf_mtime >= pptx_mtime:
            pdf_content = cached_pdf.read_bytes()
            return StreamingResponse(
                iter([pdf_content]),
                media_type="application/pdf",
                headers={
                    "Content-Disposition": "inline; filename=\"preview.pdf\"; filename*=UTF-8''"
                    + urllib.parse.quote(f"{path.stem}.pdf")
                },
            )

    import tempfile

    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            proc = await asyncio.create_subprocess_exec(
                "soffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                temp_dir,
                str(path),
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            try:
                _, _ = await asyncio.wait_for(proc.communicate(), timeout=60)
            except asyncio.TimeoutError:
                proc.kill()
                await proc.wait()
                return None
            if proc.returncode != 0:
                return None
            pdf_files = list(Path(temp_dir).glob("*.pdf"))
            if not pdf_files:
                return None
            pdf_content = pdf_files[0].read_bytes()

            # Cache the PDF next to the PPTX for future requests
            try:
                cached_pdf.write_bytes(pdf_content)
            except Exception:
                pass

            return StreamingResponse(
                iter([pdf_content]),
                media_type="application/pdf",
                headers={
                    "Content-Disposition": "inline; filename=\"preview.pdf\"; filename*=UTF-8''"
                    + urllib.parse.quote(f"{path.stem}.pdf")
                },
            )
    except Exception:
        return None


def _pptx_fallback_html(path: Path) -> HTMLResponse:
    import base64

    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu

    prs = Presentation(str(path))
    slide_width = prs.slide_width or Emu(12192000)  # default LAYOUT_WIDE
    slide_height = prs.slide_height or Emu(6858000)

    # Aspect ratio for slide container
    aspect = (
        Emu(slide_height).inches / Emu(slide_width).inches if slide_width else 0.5625
    )

    html_content = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <style>
            * {{ margin: 0; padding: 0; box-sizing: border-box; }}
            body {{ font-family: Arial, sans-serif; background: #1a1a2e; padding: 20px; }}
            .slide-container {{ max-width: 960px; margin: 16px auto; border-radius: 8px;
                overflow: hidden; box-shadow: 0 4px 20px rgba(0,0,0,0.3);
                position: relative; width: 100%; padding-top: {aspect * 100:.2f}%; }}
            .slide-inner {{ position: absolute; top: 0; left: 0; width: 100%; height: 100%;
                overflow: hidden; }}
            .slide-element {{ position: absolute; }}
            .slide-element.slide-text {{ z-index: 2; overflow: visible; }}
            .slide-element.slide-img {{ z-index: 1; overflow: hidden; }}
            .slide-element p {{ margin: 2px 0; }}
            .slide-image {{ max-width: 100%; max-height: 100%; object-fit: contain; }}
            .slide-label {{ text-align: center; color: #888; font-size: 12px;
                margin: 4px auto 12px; max-width: 960px; }}
        </style>
    </head>
    <body>
    """

    total_slides = len(prs.slides)
    sw = float(Emu(slide_width).inches) if slide_width else 13.333
    sh = float(Emu(slide_height).inches) if slide_height else 7.5

    for slide_num, slide in enumerate(prs.slides, 1):
        # Background color
        bg_color = "#FFFFFF"
        try:
            bg = slide.background
            if bg and bg.fill and bg.fill.fore_color:
                bg_color = f"#{bg.fill.fore_color.rgb}"
        except Exception:
            pass

        html_content += (
            f'<div class="slide-container">'
            f'<div class="slide-inner" style="background:{bg_color};">'
        )

        for shape in slide.shapes:
            # Position as percentages
            left_pct = (float(Emu(shape.left).inches) / sw * 100) if shape.left else 0
            top_pct = (float(Emu(shape.top).inches) / sh * 100) if shape.top else 0
            w_pct = (float(Emu(shape.width).inches) / sw * 100) if shape.width else 0
            h_pct = (float(Emu(shape.height).inches) / sh * 100) if shape.height else 0

            # Clamp width so element doesn't exceed slide right edge
            w_clamped = min(w_pct, 100 - left_pct)

            # Image shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_style = (
                    f"left:{left_pct:.1f}%;top:{top_pct:.1f}%;"
                    f"width:{w_clamped:.1f}%;height:{h_pct:.1f}%;"
                )
                try:
                    img_blob = shape.image.blob
                    img_ct = shape.image.content_type
                    b64 = base64.b64encode(img_blob).decode()
                    html_content += (
                        f'<div class="slide-element slide-img" style="{img_style}">'
                        f'<img class="slide-image" src="data:{img_ct};base64,{b64}" />'
                        f"</div>"
                    )
                except Exception:
                    pass
                continue

            # Text shapes — use auto height so text is never clipped
            text = getattr(shape, "text", None)
            if text and text.strip():
                txt_style = (
                    f"left:{left_pct:.1f}%;top:{top_pct:.1f}%;width:{w_clamped:.1f}%;"
                )
                # Only set explicit height if shape actually has one
                if h_pct > 1:
                    txt_style += f"height:{h_pct:.1f}%;"

                # Try to get font properties from the first run
                font_size = 14
                font_color = "#333333"
                font_bold = False
                try:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                font_size = int(Emu(run.font.size).pt)
                            if run.font.color and run.font.color.rgb:
                                font_color = f"#{run.font.color.rgb}"
                            if run.font.bold:
                                font_bold = True
                            break
                        break
                except Exception:
                    pass

                # Scale font size for responsive display
                scaled_size = max(8, min(font_size * 0.6, 36))
                bold_str = "font-weight:bold;" if font_bold else ""
                text_escaped = (
                    str(text)
                    .replace("&", "&amp;")
                    .replace("<", "&lt;")
                    .replace(">", "&gt;")
                    .replace("\n", "<br>")
                )
                html_content += (
                    f'<div class="slide-element slide-text" style="{txt_style}'
                    f"color:{font_color};font-size:{scaled_size:.0f}px;{bold_str}"
                    f'padding:4px;">'
                    f"<p>{text_escaped}</p></div>"
                )

        html_content += "</div></div>"
        html_content += (
            f'<div class="slide-label">Slide {slide_num} / {total_slides}</div>'
        )

    html_content += "</body></html>"
    return HTMLResponse(content=html_content)


@file_router.post("/upload")
async def upload_file(
    file: UploadFile | None = File(None),
    files: list[UploadFile] | None = File(None),
    task_type: str = Form(...),
    message: str = Form(""),
    task_id: str = Form(None),
    folder: str = Form(None),
    user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> Dict[str, Any]:
    del message
    upload_items: list[UploadFile] = []
    if file is not None:
        upload_items.append(file)
    if files:
        upload_items.extend(files)

    if not upload_items:
        raise HTTPException(status_code=422, detail="No files provided")

    single_file_mode = file is not None and (not files)
    parsed_task_id = _parse_task_id(task_id)
    uploaded_files = []

    for uploaded in upload_items:
        if not uploaded.filename or not uploaded.filename.strip():
            raise HTTPException(status_code=422, detail="No filename provided")
        if not is_allowed_file(uploaded.filename, task_type):
            raise HTTPException(
                status_code=500,
                detail=f"File type {Path(uploaded.filename).suffix.lower()} not supported for task type {task_type}",
            )

        content = await uploaded.read()
        if len(content) > MAX_FILE_SIZE:
            raise HTTPException(
                status_code=500,
                detail=f"File size exceeds maximum limit of {MAX_FILE_SIZE // (1024 * 1024)}MB",
            )

        target_path = _build_unique_file_path(
            get_upload_path(uploaded.filename, task_id, folder, _user_id_value(user))
        )
        with open(target_path, "wb") as buffer:
            buffer.write(content)

        file_record = UploadedFile(
            user_id=_user_id_value(user),
            task_id=parsed_task_id,
            filename=Path(uploaded.filename).name,
            storage_path=str(target_path),
            mime_type=uploaded.content_type,
            file_size=len(content),
        )
        db.add(file_record)
        db.flush()

        content_preview = ""
        try:
            preview_content = read_file(str(target_path))
            content_preview = (
                preview_content[:500] + "..."
                if isinstance(preview_content, str) and len(preview_content) > 500
                else preview_content
            )
        except Exception:
            content_preview = ""

        uploaded_files.append(
            {
                "file_id": file_record.file_id,
                "filename": file_record.filename,
                "file_size": file_record.file_size,
                "mime_type": file_record.mime_type,
                "content_preview": content_preview,
            }
        )

    db.commit()

    if single_file_mode:
        first_file = uploaded_files[0]
        return {
            "success": True,
            "file_id": first_file["file_id"],
            "filename": first_file["filename"],
            "file_size": first_file["file_size"],
            "mime_type": first_file["mime_type"],
            "task_type": task_type,
            "content_preview": first_file["content_preview"],
            "message": f"Successfully uploaded {first_file['filename']}",
        }

    return {
        "success": True,
        "files": uploaded_files,
        "total_files": len(uploaded_files),
        "task_type": task_type,
        "message": f"Successfully uploaded {len(uploaded_files)} files",
    }


@file_router.get("/list")
async def list_files(
    user: User = Depends(get_current_user), db: Session = Depends(get_db)
) -> Dict[str, Any]:
    query = db.query(UploadedFile)
    if not _is_admin_user(user):
        query = query.filter(UploadedFile.user_id == _user_id_value(user))

    records = query.order_by(UploadedFile.created_at.desc()).all()
    files = []
    for record in records:
        path = Path(_file_storage_path_value(record))
        record_user_id = _file_user_id_value(record)
        relative_path = _extract_relative_path(path, record_user_id)
        files.append(
            {
                "file_id": record.file_id,
                "filename": _file_name_value(record),
                "file_size": record.file_size,
                "modified_time": _to_unix_timestamp(path, record.created_at),
                "file_type": path.suffix.lower().lstrip("."),
                "relative_path": relative_path,
                "task_id": record.task_id,
                "user_id": record_user_id,
            }
        )

    return {"files": files, "total_count": len(files)}


@file_router.get("/task/{task_id}")
async def list_task_files(
    task_id: int,
    user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> Dict[str, Any]:
    """
    Get all files for a specific task.

    More efficient than /api/files/list as it filters at database level.
    Only returns files that are already registered in the database.
    """
    # Query files for this task
    query = db.query(UploadedFile).filter(UploadedFile.task_id == task_id)

    # Permission check: only show user's own files unless admin
    if not _is_admin_user(user):
        query = query.filter(UploadedFile.user_id == _user_id_value(user))

    records = query.order_by(UploadedFile.created_at.desc()).all()

    files = []
    for record in records:
        path = Path(_file_storage_path_value(record))
        if not path.exists():
            # Skip files that no longer exist on disk
            continue

        record_user_id = _file_user_id_value(record)
        relative_path = _extract_relative_path(path, record_user_id)

        # Categorize by directory (input/output/temp)
        path_parts = relative_path.split("/")
        file_category = "other"
        if len(path_parts) >= 2:
            subdir = path_parts[1]  # e.g., "input", "output", "temp"
            if subdir in ["input", "output", "temp"]:
                file_category = subdir

        files.append(
            {
                "file_id": record.file_id,
                "filename": _file_name_value(record),
                "file_size": record.file_size,
                "modified_time": _to_unix_timestamp(path, record.created_at),
                "file_type": path.suffix.lower().lstrip("."),
                "relative_path": relative_path,
                "category": file_category,
                "task_id": record.task_id,
                "user_id": record_user_id,
            }
        )

    return {"files": files, "total_count": len(files), "task_id": task_id}


@file_router.get("/download/{file_id:path}", response_model=None)
async def download_file(
    file_id: str,
    user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> Any:
    file_record, full_path, owner_user_id = _resolve_file_path(
        db, file_id, _user_id_value(user)
    )

    # Check access permissions
    if file_record:
        _check_file_access(file_record, user)
        file_name = _file_name_value(file_record)
        media_type = _guess_media_type(file_name)
    else:
        # For legacy files without records, check ownership
        if owner_user_id != _user_id_value(user) and not _is_admin_user(user):
            raise HTTPException(status_code=403, detail="Access denied")
        file_name = full_path.name
        media_type = _guess_media_type(file_name)

    _ensure_under_uploads(full_path, owner_user_id)

    if not full_path.exists():
        raise HTTPException(status_code=404, detail="File not found")

    converted_pdf = await _try_convert_pptx_to_pdf(full_path)
    if converted_pdf is not None:
        return converted_pdf

    # For images and other viewable content, set Content-Disposition to inline
    # to allow browser to display the file instead of downloading it
    content_disposition = (
        "inline"
        if media_type.startswith(("image/", "video/", "audio/", "text/"))
        else "attachment"
    )

    return FileResponse(
        path=str(full_path),
        filename=file_name,
        media_type=media_type,
        headers={
            "Content-Disposition": f'{content_disposition}; filename="{file_name}"'
        },
    )


@file_router.get("/preview/{file_id:path}", response_model=None)
async def preview_file(
    file_id: str,
    user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> Any:
    file_record, full_path, owner_user_id = _resolve_file_path(
        db, file_id, _user_id_value(user)
    )

    # Check access permissions
    if file_record:
        _check_file_access(file_record, user)
        file_name = _file_name_value(file_record)
        media_type = _guess_media_type(file_name)
    else:
        # For legacy files without records, check ownership
        if owner_user_id != _user_id_value(user) and not _is_admin_user(user):
            raise HTTPException(status_code=403, detail="Access denied")
        file_name = full_path.name
        media_type = _guess_media_type(file_name)

    _ensure_under_uploads(full_path, owner_user_id)

    if not full_path.exists():
        raise HTTPException(status_code=404, detail="File not found")

    # PPTX preview: use HTML (preserves Chinese text + images)
    # PDF conversion loses Chinese text due to font mapping issues in LibreOffice
    if full_path.suffix.lower() == ".pptx":
        try:
            return _pptx_fallback_html(full_path)
        except Exception:
            pass

    return FileResponse(
        path=str(full_path),
        filename=file_name,
        media_type=media_type,
        headers={"Content-Disposition": "inline"},
    )


@file_router.get("/public/preview/{file_id:path}", response_model=None)
async def public_preview_file(
    file_id: str,
    relative_path: Optional[str] = Query(default=None),
    db: Session = Depends(get_db),
) -> Any:
    # For public preview, we need to handle both file_id and legacy paths
    # Try UUID first
    file_record = None
    base_path = None
    owner_user_id = None

    if is_valid_uuid(file_id):
        file_record = (
            db.query(UploadedFile).filter(UploadedFile.file_id == file_id).first()
        )

    if file_record:
        base_path = Path(_file_storage_path_value(file_record))
        owner_user_id = _file_user_id_value(file_record)
    else:
        # Try to resolve as legacy path across all user directories
        result = resolve_legacy_file_path_cross_user(file_id)
        if result is None:
            raise HTTPException(status_code=404, detail="File not found")

        base_path, owner_user_id = result

    target_path = _resolve_public_preview_target(
        base_path,
        relative_path,
        owner_user_id,
    )

    if not target_path.exists() or not target_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")

    converted_pdf = await _try_convert_pptx_to_pdf(target_path)
    if converted_pdf is not None:
        return converted_pdf

    return FileResponse(
        path=str(target_path),
        filename=target_path.name,
        media_type=_guess_media_type(target_path.name),
        headers={"Content-Disposition": "inline"},
    )


@file_router.post("/backfill")
async def backfill_files(
    user: User = Depends(get_current_user), db: Session = Depends(get_db)
) -> Dict[str, Any]:
    """
    Manually trigger file backfill to sync filesystem with database.

    This is a maintenance operation that scans the filesystem and creates
    database records for any unregistered files. Only available to admins.
    """
    if not _is_admin_user(user):
        raise HTTPException(status_code=403, detail="Admin access required")

    try:
        _backfill_uploaded_file_records(db, user)
        return {"success": True, "message": "File backfill completed successfully"}
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Backfill failed: {str(e)}") from e


@file_router.delete("/{file_id:path}")
async def delete_file(
    file_id: str,
    user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> Dict[str, Any]:
    file_record, file_path, owner_user_id = _resolve_file_path(
        db, file_id, _user_id_value(user)
    )

    # Check access permissions
    if file_record:
        _check_file_access(file_record, user)
        file_name = _file_name_value(file_record)
    else:
        # For legacy files without records, check ownership
        if owner_user_id != _user_id_value(user) and not _is_admin_user(user):
            raise HTTPException(status_code=403, detail="Access denied")
        file_name = file_path.name

    _ensure_under_uploads(file_path, owner_user_id)

    if file_path.exists() and file_path.is_file():
        file_path.unlink()

    # Delete database record if exists
    if file_record:
        db.delete(file_record)
        db.commit()

    return {
        "success": True,
        "message": f"File {file_name} deleted successfully",
        "file_id": file_id,
    }
